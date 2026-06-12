from pathlib import Path

from pptx2web.metadata import extract


def test_deck_basics(deck_path: Path):
    deck = extract(deck_path)
    assert deck.title == "Deck de prueba"
    assert deck.slide_count == 5
    assert deck.slide_width_pt > 0 and deck.slide_height_pt > 0


def test_titles(deck_path: Path):
    deck = extract(deck_path)
    assert [s.title for s in deck.slides] == [f"Título {i}" for i in range(1, 6)]


def test_search_text(deck_path: Path):
    deck = extract(deck_path)
    assert "único-3" in deck.slides[2].search_text
    assert "Título 3" in deck.slides[2].search_text


def test_notes_html_escaped(deck_path: Path):
    deck = extract(deck_path)
    notes1 = deck.slides[0].notes_html
    assert notes1 is not None
    assert "<p>Nota del orador 1</p>" in notes1
    assert "&lt;especial&gt;" in notes1  # HTML escapado
    assert deck.slides[1].notes_html is None  # slide par: sin notas


def test_transitions(deck_path: Path):
    deck = extract(deck_path)
    # slide 1: sin <p:transition> → cut/0
    assert deck.slides[0].transition.type == "cut"
    assert deck.slides[0].transition.duration == 0
    # slide 2: fade con dur=700
    assert deck.slides[1].transition.type == "fade"
    assert deck.slides[1].transition.duration == 700
    # slide 3: push, dur por defecto 500
    assert deck.slides[2].transition.type == "push"
    assert deck.slides[2].transition.duration == 500
    # slides 4 y 5: wipe y split
    assert deck.slides[3].transition.type == "wipe"
    assert deck.slides[4].transition.type == "split"


def test_morph_degrades_to_fade(tmp_path: Path):
    from conftest import build_deck

    deck = extract(build_deck(tmp_path / "m.pptx", n_slides=7))
    # slide 7 recibe 'morph' (índice 5 del ciclo) → degrada a fade
    assert deck.slides[6].transition.type == "fade"


def test_no_media(deck_path: Path):
    deck = extract(deck_path)
    assert all(s.media == [] for s in deck.slides)


def test_text_run_link_extracted(deck_path: Path):
    deck = extract(deck_path)
    links = deck.slides[1].links  # slide 2
    ext = [l for l in links if l.href == "https://example.com/video"]
    assert len(ext) == 1
    assert ext[0].kind == "text"
    assert ext[0].tooltip == "Ver video"
    r = ext[0].rect
    assert 0 <= r.x <= 1 and 0 < r.w <= 1 and 0 < r.h <= 1


def test_shape_link_extracted(deck_path: Path):
    deck = extract(deck_path)
    links = deck.slides[2].links  # slide 3
    ext = [l for l in links if l.href == "https://example.com/pagina"]
    assert len(ext) == 1
    assert ext[0].kind == "shape"
    assert ext[0].slide is None


def test_internal_link_extracted(deck_path: Path):
    deck = extract(deck_path)
    links = deck.slides[3].links  # slide 4
    internal = [l for l in links if l.slide is not None]
    assert len(internal) == 1
    assert internal[0].slide == 1
    assert internal[0].href is None


def test_quiz_parsed_from_notes(deck_path: Path):
    deck = extract(deck_path)
    s5 = deck.slides[4]
    assert s5.quiz is not None
    assert s5.quiz.question == "¿Cuál es la 4a revolución industrial?"
    assert [o.text for o in s5.quiz.options] == [
        "Mecanización", "Digitalización", "Electricidad"
    ]
    assert [o.correct for o in s5.quiz.options] == [False, True, False]
    assert s5.quiz.feedback_ok == "¡Correcto!"
    assert s5.quiz.feedback_ko == "Revisa la lámina 2."
    # las notas [quiz] no se publican como notas
    assert s5.notes_html is None
    assert deck.warnings == []


def test_malformed_quiz_warns_and_keeps_notes(tmp_path: Path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "T"
    slide.notes_slide.notes_text_frame.text = "[quiz]\nSolo una\n+ única opción"
    p = tmp_path / "bad.pptx"
    prs.save(str(p))

    deck = extract(p)
    assert deck.slides[0].quiz is None
    assert deck.slides[0].notes_html is not None  # quedan como notas normales
    assert any("malformado" in w for w in deck.warnings)


def test_quiz_requires_single_correct():
    from pptx2web.metadata import _parse_quiz

    quiz, err = _parse_quiz("[quiz]\n+ a\n+ b")
    assert quiz is None and "1 opción correcta" in err
    quiz, err = _parse_quiz("[quiz]\n- a\n- b")
    assert quiz is None and "1 opción correcta" in err
    quiz, err = _parse_quiz("[quiz]\n- a\n+ b")
    assert quiz is not None and quiz.question is None
