"""Fixtures compartidos: genera decks .pptx de prueba con python-pptx.

Ninguno de estos tests requiere COM/PowerPoint; corren en CI.
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "src"))

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from lxml import etree  # noqa: E402

NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _add_transition(slide, kind: str, duration_ms: int | None = None) -> None:
    """Inyecta <p:transition> en el XML del slide (python-pptx no tiene API)."""
    el = etree.SubElement(slide._element, f"{{{NS_P}}}transition")
    if duration_ms is not None:
        el.set(f"{{{NS_P}}}dur", str(duration_ms))
    etree.SubElement(el, f"{{{NS_P}}}{kind}")


def build_deck(path: Path, n_slides: int = 3) -> Path:
    prs = Presentation()
    prs.core_properties.title = "Deck de prueba"
    layout = prs.slide_layouts[1]  # título + contenido
    transitions = ["fade", "push", "wipe", "split", "cut", "morph"]

    for i in range(1, n_slides + 1):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Título {i}"
        body = slide.placeholders[1]
        body.text_frame.text = f"Contenido del slide {i} con texto buscable único-{i}"

        if i == 5:
            # slide de quiz: las notas llevan el bloque [quiz]
            slide.notes_slide.notes_text_frame.text = (
                "[quiz]\n¿Cuál es la 4a revolución industrial?\n"
                "- Mecanización\n+ Digitalización\n- Electricidad\n"
                "ok: ¡Correcto!\nno: Revisa la lámina 2."
            )
        elif i % 2 == 1:
            notes = slide.notes_slide.notes_text_frame
            notes.text = f"Nota del orador {i}\nSegunda línea & <especial>"

        if i == 2:
            # link externo en un run de texto
            p = body.text_frame.add_paragraph()
            run = p.add_run()
            run.text = "Ver video"
            run.hyperlink.address = "https://example.com/video"
        if i == 3:
            # link externo sobre un shape completo
            box = slide.shapes.add_textbox(
                Inches(1), Inches(5), Inches(3), Inches(1)
            )
            box.text_frame.text = "Botón externo"
            box.click_action.hyperlink.address = "https://example.com/pagina"
        if i == 4:
            # link interno a la primera lámina
            box = slide.shapes.add_textbox(
                Inches(1), Inches(5), Inches(3), Inches(1)
            )
            box.text_frame.text = "Volver al inicio"
            box.click_action.target_slide = prs.slides[0]

        if i > 1:  # el primero sin transición
            _add_transition(slide, transitions[(i - 2) % len(transitions)],
                            duration_ms=700 if i == 2 else None)

    prs.save(str(path))
    return path


@pytest.fixture
def deck_path(tmp_path: Path) -> Path:
    return build_deck(tmp_path / "fixture.pptx", n_slides=5)
