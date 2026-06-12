"""Extracción de metadatos con python-pptx (sin COM, testeable en CI)."""
from __future__ import annotations

import html
import logging
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from .models import (
    DEFAULT_TRANSITION_DURATION_MS,
    Deck,
    LinkItem,
    MediaItem,
    Quiz,
    QuizOption,
    Rect,
    SlideMeta,
    Transition,
)

log = logging.getLogger("pptx2web")

# python-pptx no expone transiciones con API de alto nivel; se lee el XML
# del slide. Namespace de PresentationML:
NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Elementos hijos de <p:transition> que mapean directo a tipos del player.
# Cualquier otro (morph, zoom, prism…) degrada a 'fade'.
_DIRECT_TRANSITIONS = {"fade", "push", "wipe", "split", "cut"}


def extract(pptx_path: Path) -> Deck:
    prs = Presentation(str(pptx_path))

    title = (prs.core_properties.title or "").strip() or pptx_path.stem
    warnings: list[str] = []

    # Mapas para resolver links internos (a otra lámina): tanto la API de
    # alto nivel (target_slide) como las rels XML (partname) → orden 1-based.
    id_to_index = {s.slide_id: i for i, s in enumerate(prs.slides, start=1)}
    part_to_index = {str(s.part.partname): i for i, s in enumerate(prs.slides, start=1)}

    slides: list[SlideMeta] = []
    for i, slide in enumerate(prs.slides, start=1):
        quiz, notes_html = _quiz_or_notes(slide, i, warnings)
        slides.append(
            SlideMeta(
                index=i,
                title=_slide_title(slide, i),
                notes_html=notes_html,
                search_text=_search_text(slide),
                transition=_transition(slide),
                media=_media_items(slide, prs.slide_width, prs.slide_height),
                links=_link_items(
                    slide, prs.slide_width, prs.slide_height,
                    id_to_index, part_to_index, warnings,
                ),
                quiz=quiz,
            )
        )

    return Deck(
        title=title,
        slide_width_pt=prs.slide_width / 12700,
        slide_height_pt=prs.slide_height / 12700,
        slides=slides,
        warnings=warnings,
    )


def _slide_title(slide, index: int) -> str:
    try:
        ph = slide.shapes.title
        if ph is not None and ph.has_text_frame:
            text = ph.text_frame.text.strip()
            if text:
                return _first_line(text)
    except Exception:
        pass
    # Sin placeholder de título: primera línea de texto de cualquier shape
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return _first_line(text)
    return f"Diapositiva {index}"


def _first_line(text: str, max_len: int = 120) -> str:
    line = text.splitlines()[0].strip()
    return line[:max_len]


def _notes_html(slide) -> str | None:
    if not slide.has_notes_slide:
        return None
    tf = slide.notes_slide.notes_text_frame
    if tf is None:
        return None
    paragraphs = [p.text.strip() for p in tf.paragraphs]
    paragraphs = [p for p in paragraphs if p]
    if not paragraphs:
        return None
    return "".join(f"<p>{html.escape(p)}</p>" for p in paragraphs)


def _search_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        chunks.extend(_shape_text(shape))
    return " ".join(" ".join(c.split()) for c in chunks if c.strip())


def _shape_text(shape) -> list[str]:
    chunks: list[str] = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            chunks.extend(_shape_text(sub))
        return chunks
    if shape.has_text_frame:
        text = shape.text_frame.text
        if text:
            chunks.append(text)
    if getattr(shape, "has_table", False) and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text:
                    chunks.append(cell.text)
    return chunks


def _transition(slide) -> Transition:
    el = slide._element.find(f"{{{NS_P}}}transition")
    if el is None:
        return Transition(type="cut", duration=0)

    duration = DEFAULT_TRANSITION_DURATION_MS
    dur_attr = el.get(f"{{{NS_P}}}dur") or el.get("dur")
    if dur_attr and dur_attr.isdigit():
        duration = int(dur_attr)

    ttype = "fade"
    for child in el:
        local = child.tag.rsplit("}", 1)[-1]
        if local in _DIRECT_TRANSITIONS:
            ttype = local
            break
    if ttype == "cut":
        duration = 0
    return Transition(type=ttype, duration=duration)


def _quiz_or_notes(slide, index: int, warnings: list[str]):
    """Si las notas empiezan con [quiz] devuelve (Quiz, None); si no,
    (None, notes_html). Un quiz malformado deja las notas visibles y avisa."""
    text = _notes_text(slide)
    if text is None or not text.lstrip().lower().startswith("[quiz]"):
        return None, _notes_html(slide)
    quiz, error = _parse_quiz(text)
    if quiz is None:
        warnings.append(f"Slide {index}: bloque [quiz] malformado ({error}); "
                        "se publica como notas normales")
        return None, _notes_html(slide)
    return quiz, None


def _notes_text(slide) -> str | None:
    if not slide.has_notes_slide:
        return None
    tf = slide.notes_slide.notes_text_frame
    return tf.text if tf is not None else None


def _parse_quiz(text: str) -> tuple[Quiz | None, str | None]:
    """Sintaxis (una cosa por línea, tras la marca [quiz]):

        [quiz]
        ¿Pregunta opcional?
        - opción incorrecta
        + opción correcta
        ok: feedback al acertar (opcional)
        no: feedback al fallar (opcional)
    """
    lines = text.strip().splitlines()
    lines[0] = lines[0].strip()[len("[quiz]"):]  # permitir contenido tras la marca

    question_lines: list[str] = []
    options: list[QuizOption] = []
    feedback_ok: str | None = None
    feedback_ko: str | None = None

    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if line.startswith("-"):
            options.append(QuizOption(text=line[1:].strip(), correct=False))
        elif line.startswith("+"):
            options.append(QuizOption(text=line[1:].strip(), correct=True))
        elif line.lower().startswith("ok:"):
            feedback_ok = line[3:].strip() or None
        elif line.lower().startswith("no:"):
            feedback_ko = line[3:].strip() or None
        elif options:
            return None, f"línea inesperada después de las opciones: '{line[:40]}'"
        else:
            question_lines.append(line)

    if len(options) < 2:
        return None, "se requieren al menos 2 opciones (líneas con - o +)"
    correct_count = sum(1 for o in options if o.correct)
    if correct_count != 1:
        return None, f"debe haber exactamente 1 opción correcta (+), hay {correct_count}"
    if any(not o.text for o in options):
        return None, "hay opciones sin texto"

    question = " ".join(question_lines).strip() or None
    return Quiz(question=question, options=options,
                feedback_ok=feedback_ok, feedback_ko=feedback_ko), None


_PPACTION_JUMP = "ppaction://hlinksldjump"


def _link_items(
    slide,
    slide_w_emu: int,
    slide_h_emu: int,
    id_to_index: dict,
    part_to_index: dict,
    warnings: list[str],
) -> list[LinkItem]:
    items: list[LinkItem] = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Las coordenadas de los hijos son locales al grupo: si algún
            # descendiente tiene link, el hotspot es el bbox del grupo entero
            # (aproximación documentada).
            group_links = _collect_group_links(shape, slide, id_to_index,
                                               part_to_index, warnings)
            rect = _shape_rect(shape, slide_w_emu, slide_h_emu)
            if rect:
                for kind, href, target, tooltip in group_links:
                    items.append(LinkItem(kind=kind, rect=rect, href=href,
                                          slide=target, tooltip=tooltip))
            continue
        items.extend(
            _shape_links(shape, slide, slide_w_emu, slide_h_emu,
                         id_to_index, part_to_index, warnings)
        )
    return items


def _shape_links(
    shape, slide, slide_w_emu, slide_h_emu, id_to_index, part_to_index, warnings
) -> list[LinkItem]:
    rect = _shape_rect(shape, slide_w_emu, slide_h_emu)
    if rect is None:
        return []
    found = _collect_shape_links(shape, slide, id_to_index, part_to_index, warnings)
    items = [
        LinkItem(kind=kind, rect=rect, href=href, slide=target, tooltip=tooltip)
        for kind, href, target, tooltip in found
    ]
    text_links = [i for i in items if i.kind == "text"]
    if len(text_links) > 1:
        warnings.append(
            f"Slide {slide_index_of(slide, id_to_index)}: la caja '{shape.name}' "
            f"tiene {len(text_links)} links de texto distintos; las zonas "
            "clicables quedan superpuestas (ocupan toda la caja)"
        )
    return items


def slide_index_of(slide, id_to_index: dict) -> int:
    return id_to_index.get(slide.slide_id, 0)


def _collect_group_links(group, slide, id_to_index, part_to_index, warnings):
    found = []
    for sub in group.shapes:
        if sub.shape_type == MSO_SHAPE_TYPE.GROUP:
            found.extend(_collect_group_links(sub, slide, id_to_index,
                                              part_to_index, warnings))
        else:
            found.extend(_collect_shape_links(sub, slide, id_to_index,
                                              part_to_index, warnings))
    return found


def _collect_shape_links(shape, slide, id_to_index, part_to_index, warnings):
    """Devuelve tuplas (kind, href, slide_target, tooltip) sin rect."""
    found: list[tuple] = []
    seen: set[tuple] = set()

    def add(kind, href, target, tooltip):
        key = (href, target)
        if (href or target) and key not in seen:
            seen.add(key)
            found.append((kind, href, target, tooltip))

    # 1) Link del shape completo (click_action)
    try:
        action = shape.click_action
    except (AttributeError, ValueError):
        action = None
    if action is not None:
        try:
            target_slide = action.target_slide
        except (KeyError, ValueError):
            target_slide = None
        if target_slide is not None:
            add("shape", None, id_to_index.get(target_slide.slide_id),
                shape.name or None)
        elif action.hyperlink.address:
            add("shape", action.hyperlink.address, None, shape.name or None)

    # 2) Links en runs de texto (rect = caja completa, kind="text")
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                addr = run.hyperlink.address
                if addr:
                    add("text", addr, None, run.text.strip() or None)
                    continue
                # interno: <a:hlinkClick action="ppaction://hlinksldjump" r:id=…>
                rpr = run._r.find(f"{{{ns_a}}}rPr")
                if rpr is None:
                    continue
                hl = rpr.find(f"{{{ns_a}}}hlinkClick")
                if hl is None or hl.get("action") != _PPACTION_JUMP:
                    continue
                rid = hl.get(f"{{{ns_r}}}id")
                if not rid:
                    continue
                try:
                    partname = str(slide.part.rels[rid].target_part.partname)
                except (KeyError, AttributeError):
                    continue
                target = part_to_index.get(partname)
                if target:
                    add("text", None, target, run.text.strip() or None)
    return found


def _shape_rect(shape, slide_w_emu: int, slide_h_emu: int) -> Rect | None:
    try:
        rect = Rect(
            x=shape.left / slide_w_emu,
            y=shape.top / slide_h_emu,
            w=shape.width / slide_w_emu,
            h=shape.height / slide_h_emu,
        )
    except TypeError:
        return None  # shape sin posición (raro: placeholders heredados)
    if rect.w <= 0 or rect.h <= 0:
        return None
    return rect


def _media_items(slide, slide_w_emu: int, slide_h_emu: int) -> list[MediaItem]:
    items: list[MediaItem] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.MEDIA:
            continue
        kind = _media_kind(shape)
        if kind is None:
            continue
        part_name = _media_part_name(shape)
        if part_name is None:
            log.warning(
                "Slide %s: shape de media sin part asociado, se omite",
                slide.slide_id,
            )
            continue
        try:
            rect = Rect(
                x=Emu(shape.left) / slide_w_emu,
                y=Emu(shape.top) / slide_h_emu,
                w=Emu(shape.width) / slide_w_emu,
                h=Emu(shape.height) / slide_h_emu,
            )
        except TypeError:
            rect = Rect(x=0.1, y=0.1, w=0.8, h=0.8)
        items.append(MediaItem(type=kind, source_part=part_name, rect=rect))
    return items


_VIDEO_EXTS = {".mp4", ".m4v", ".mov", ".wmv", ".avi", ".mpg", ".mpeg", ".asf", ".webm"}
_AUDIO_EXTS = {".mp3", ".wav", ".wma", ".m4a", ".aac", ".ogg"}


def _media_kind(shape) -> str | None:
    name = _media_part_name(shape)
    if name:
        ext = Path(name).suffix.lower()
        if ext in _VIDEO_EXTS:
            return "video"
        if ext in _AUDIO_EXTS:
            return "audio"
    return None


def _media_part_name(shape) -> str | None:
    """Resuelve el part de ppt/media/ del shape vía la relación r:link/r:embed
    del elemento <a:videoFile>/<a:audioFile>."""
    ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    el = shape._element
    for tag in ("videoFile", "audioFile", "quickTimeFile"):
        media_el = el.find(f".//{{{ns_a}}}{tag}")
        if media_el is None:
            continue
        rid = media_el.get(f"{{{ns_r}}}link") or media_el.get(f"{{{ns_r}}}embed")
        if not rid:
            continue
        try:
            part = shape.part.rels[rid].target_part
            return part.partname.split("/")[-1]
        except (KeyError, AttributeError):
            # Relación externa (link a archivo fuera del .pptx): no empaquetable
            return None
    return None
